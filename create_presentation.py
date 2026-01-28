from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 121)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    # Add accent line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
    line.line.color.rgb = RGBColor(31, 78, 121)
    line.line.width = Pt(3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content_items:
        if isinstance(item, tuple):
            # (bold_text, regular_text) format
            p = text_frame.add_paragraph()
            p.level = 0
            p.space_before = Pt(8)
            
            run = p.add_run()
            run.text = item[0]
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(31, 78, 121)
            
            run = p.add_run()
            run.text = " " + item[1]
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(68, 68, 68)
        else:
            # Regular bullet point
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(68, 68, 68)
            p.space_before = Pt(10)
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    # Add accent line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
    line.line.color.rgb = RGBColor(31, 78, 121)
    line.line.width = Pt(3)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.4), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for item in left_items:
        if isinstance(item, tuple):
            p = left_frame.add_paragraph()
            p.level = 0
            run = p.add_run()
            run.text = item[0]
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(31, 78, 121)
            run = p.add_run()
            run.text = " " + item[1]
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(68, 68, 68)
        else:
            p = left_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(68, 68, 68)
        p.space_before = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.1), Inches(1.6), Inches(4.4), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for item in right_items:
        if isinstance(item, tuple):
            p = right_frame.add_paragraph()
            p.level = 0
            run = p.add_run()
            run.text = item[0]
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(31, 78, 121)
            run = p.add_run()
            run.text = " " + item[1]
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(68, 68, 68)
        else:
            p = right_frame.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(68, 68, 68)
        p.space_before = Pt(6)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "Sales Dashboard", "Interactive Retail Performance Analytics\n2014-2017 Analysis")

# Slide 2: Executive Overview
add_content_slide(prs, "Executive Overview", [
    "Analyzes 9,995+ sales transactions (2014-2017)",
    "Interactive dashboards for sales, profit, and year-over-year analysis",
    "Dynamic filtering by year enables trend identification",
    "Integrates sales, manager, and return data for complete visibility",
    ("Business Value:", "Transform raw transaction data into actionable insights for regional and category-level decision making.")
])

# Slide 3: Sales Analysis Dashboard
add_content_slide(prs, "Sales Analysis Dashboard", [
    ("Regional Performance Breakdown:", "Sales by region with segment drill-down capability"),
    "Identify high-performing and underperforming regions",
    "Track sales velocity by product category",
    ("Segment Analysis:", "Consumer, Corporate, and Home Office performance comparison"),
    "Category preferences by customer segment",
    ("Category Deep-Dive:", "Top and bottom performing product categories"),
    "[Screenshot placeholder: Sales Analysis Dashboard visual]"
])

# Slide 4: Profit Analysis Dashboard
add_content_slide(prs, "Profit Analysis Dashboard", [
    ("Margin Analysis:", "Profit by region and category"),
    "Profitability trends across 2014-2017",
    "Identify margin-erosion categories",
    ("Discount Impact Assessment:", "Correlate discount depth with profit levels"),
    "Identify high-discount, low-profit scenarios",
    ("Performance Ranking:", "Regional profitability hierarchy"),
    "Profit contribution % by segment",
    ("Key Insight:", "Profit trends reveal which categories are truly valuable vs. high-volume but low-margin.")
])

# Slide 5: Year-Over-Year Analysis
add_content_slide(prs, "Year-Over-Year Analysis Dashboard", [
    ("Growth Metrics:", "Year-over-year sales and profit growth rates"),
    "Identify accelerating and decelerating categories",
    ("Dynamic Filtering:", "Select year parameter (2014-2017) to compare periods"),
    "Quick pivot to different time windows",
    ("Trend Patterns:", "Consistent growth categories vs. volatile ones"),
    "Seasonal patterns year-over-year",
    "Early warning indicators for declining trends"
])

# Slide 6: Data Integration
add_two_column_slide(prs, "Data Integration & Intelligence", [
    "Three Data Streams:",
    ("Orders Table:", "9,995 transactions"),
    ("People Table:", "Manager-to-region mapping"),
    ("Returns Table:", "297 returned orders"),
    "Key Dimensions:",
    "Geography: Region, State, City",
    "Business: Segment, Category",
    "Temporal: 2014-2017"
], [
    "Derived Insights:",
    ("Return Rates:", "Manager/region accountability"),
    ("Discount Patterns:", "Margin protection"),
    ("Customer Segment:", "Profitability analysis"),
    "Linked Analysis:",
    "Cross-functional visibility",
    "Accountability tracking",
    "Risk identification"
])

# Slide 7: Interactive Capabilities
add_content_slide(prs, "Interactive Capabilities", [
    ("Year-Based Filter:", "Parameter selector (2014, 2015, 2016, 2017)"),
    "All dashboard views update dynamically",
    "Enables trend comparison and deep-dive analysis",
    ("Drill-Down Capabilities:", "Region → Category → Sub-Category"),
    "Customer Segment → Product Detail",
    ("Cross-Dashboard Navigation:", "Filter selections carry across dashboards"),
    "Return patterns linked to manager performance"
])

# Slide 8: Live Dashboard Access
add_content_slide(prs, "Live Dashboard Access", [
    "Access the interactive dashboards:",
    ("Profit Analysis Dashboard:", "https://public.tableau.com/views/salesdashboard-ProfitAnalysis/ProfitAnalysis"),
    ("Year Over Year Analysis:", "https://public.tableau.com/views/salesdashboard-YearOverYearAnalysis/YearOverYearDashboard"),
    "Deployed on Tableau Public",
    "Real-time interactive exploration",
    "Dynamic filtering and drill-down capabilities"
])

# Slide 9: Key Takeaways
add_content_slide(prs, "Key Takeaways", [
    ("Problem Solved:", "Retail leadership now has real-time, consolidated visibility into sales performance"),
    ("Capability Delivered:", "Dynamic filtering, multi-dimensional drill-down, accountability tracking"),
    ("Business Impact:", "Faster decision-making (minutes vs. weeks)"),
    "Proactive identification of risks and opportunities",
    "Data-driven resource allocation",
    "Improved margin management through discount tracking"
])

# Save presentation
output_path = "/Users/rishabhkimothi/Documents/My Tableau Repository/Workbooks/sales-dashboard-github/Sales_Dashboard_Presentation.pptx"
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
